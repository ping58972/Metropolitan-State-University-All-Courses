'''Author: Nalongsone Danddank, Assignment #3, Jan 26 2022, ICS499 METRO STATE '''
import string
import requests
import json
import random
from pptx import Presentation
from pptx.util import Inches


def getCharsFromAPI(string, language):
    '''@string: String for api params 
    @language: String for api params '''
    URL_char = "https://indic-wp.thisisjava.com/api/getLogicalChars.php"
    URL_len = "https://indic-wp.thisisjava.com/api/getLength.php"
    params = {"string": string, "language": language}
    headers = {
        'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_10_1) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/39.0.2171.95 Safari/537.36'}
    # get text or string that request from API
    r_chars = requests.get(url=URL_char, params=params,
                           headers=headers, timeout=5)
    # handle text like json format from api by encode and decode the string.
    text = r_chars.text
    if text.startswith(u'\ufeff'):
        text = text.encode('utf8')[6:].decode('utf8')
    # convert text to dictionary by json
    data = json.loads(text)

    # get chars length from API
    r_len = requests.get(url=URL_len, params=params,
                         headers=headers, timeout=5)
    text = r_len.text
    if text.startswith(u'\ufeff'):
        text = text.encode('utf8')[6:].decode('utf8')
    len_json = json.loads(text)
    # return both chars list and its length
    return data["data"], len_json["data"]


def generateChars15Chunk(chars_len, chars):
    '''Generating chars list 15 chunk by random.
    @chars_len: int length of chars list that generate by api.
    @chars: list or array which generate by api'''
    chars_15chunks = []
    if chars_len > 15:
        # combine the rest of chars to the most end of char in list
        # when the length more than 15
        d, m = divmod(chars_len, 15)
        logical_list = [d for _ in range(15 - m)] + [d+1 for _ in range(m)]
        index = 0
        for value in logical_list:
            temp_char = ""
            for _ in range(value):
                temp_char += chars[index]
                index += 1
            chars_15chunks.append(temp_char)
    elif chars_len == 15:
        chars_15chunks = chars
    else:
        # when less than 15, make the dupicate the last char until 15.
        extend = 15 - chars_len
        chars_15chunks = chars + [chars[-1] for _ in range(extend)]
    # shuffle the list to make a random in list
    random.Random(4).shuffle(chars_15chunks)
    return chars_15chunks


def buildPPT(chars_15chunks):
    '''to build Slide 15 Puzzle PPT by 4X4 table
    @chars_15chunks: list of random chars'''
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    # create a size for table by 4X4
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(3)
    shape = slide.shapes.add_table(4, 4, x, y, cx, cy)
    title.text = "Pluzz Plus"
    table = shape.table
    for i, chars in enumerate(chars_15chunks):
        # put all chars from list to each cell of table
        cell = table.cell(i//4, i % 4)
        cell.text = chars
    prs.save('test.pptx')


if __name__ == "__main__":
    # string = "అమెరికాఆస్ట్రేలియా"
    # lang = "Telugu"
    # string = "A rose by any other name would smell as sweet."
    # lang = "English"
    string = "ABCDEFGHIJKLMNOPQ"
    lang = "English"
    chars, chars_len = getCharsFromAPI(string, lang)
    chars_15chunks = generateChars15Chunk(chars_len, chars)
    buildPPT(chars_15chunks)
