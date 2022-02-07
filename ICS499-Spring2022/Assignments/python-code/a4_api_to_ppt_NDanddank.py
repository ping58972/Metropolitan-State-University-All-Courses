'''Author: Nalongsone Danddank, Assignment #4, Feb 2 2022, ICS499 METRO STATE '''

import requests
import json
import random
import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from csv import reader


def getCharsFromAPI(string, language):
    '''@string: String for api params 
    @language: String for api params '''
    URL_char = "https://indic-wp.thisisjava.com/api/getLogicalChars.php"
    URL_len = "https://indic-wp.thisisjava.com/api/getLength.php"
    params = {"string": string, "language": language}
    headers = {
        'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_10_1) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/39.0.2171.95 Safari/537.36'}
    # get text or string that request from API
    r_chars = requests.get(url=URL_char, params=params,
                           headers=headers, timeout=5)
    # handle text like json format from api by encode and decode the string.
    text = r_chars.text
    if text.startswith(u'\ufeff'):
        text = text.encode('utf8')[6:].decode('utf8')
    # convert text to dictionary by json
    data = json.loads(text)

    # get chars length from API
    r_len = requests.get(url=URL_len, params=params,
                         headers=headers, timeout=5)
    text = r_len.text
    if text.startswith(u'\ufeff'):
        text = text.encode('utf8')[6:].decode('utf8')
    len_json = json.loads(text)
    # return both chars list and its length
    return data["data"], len_json["data"]


def generateChars15Chunk(chars_len, chars):
    '''Generating chars list 15 chunk by random.
    @chars_len: int length of chars list that generate by api.
    @chars: list or array which generate by api'''
    chars_15chunks = []
    if chars_len > 15:
        # combine the rest of chars to the most end of char in list
        # when the length more than 15
        d, m = divmod(chars_len, 15)
        logical_list = [d for _ in range(15 - m)] + [d+1 for _ in range(m)]
        index = 0
        for value in logical_list:
            temp_char = ""
            for _ in range(value):
                temp_char += chars[index]
                index += 1
            chars_15chunks.append(temp_char)
    elif chars_len == 15:
        chars_15chunks = chars
    else:
        # when less than 15, make the dupicate the last char until 15.
        extend = 15 - chars_len
        chars_15chunks = chars + [chars[-1] for _ in range(extend)]
    # shuffle the list to make a random in list
    random.Random(4).shuffle(chars_15chunks)
    return chars_15chunks


def addSlide(chars_15chunks, prs, slide_title):
    '''to build Slide 15 Puzzle PPT by 4X4 table
    @chars_15chunks: list of random chars'''
    # innter function to return the cells of a table
    def iter_cells(table):
        for row in table.rows:
            for cell in row.cells:
                yield cell
    # constant for the slide layout
    # https://python-pptx.readthedocs.io/en/latest/user/table.html
    SLD_LAYOUT_TITLE_AND_CONTENT = 5
    slide = prs.slides.add_slide(prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT])
    title = slide.shapes.title
    # create a size for table by 4X4
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(3)
    shape = slide.shapes.add_table(4, 4, x, y, cx, cy)
    title.text = slide_title
    # set the table style
    # https://github.com/mbachtell/python-pptx-table-styles
    style_id = '{5940675A-B579-460E-94D1-54222C63F5DA}'
    tbl = shape._element.graphic.graphicData.tbl
    tbl[0][-1].text = style_id
    table = shape.table
    for i, chars in enumerate(chars_15chunks):
        # put all chars from list to each cell of table
        cell = table.cell(i//4, i % 4)
        cell.text = chars
    # set the font for the table cells
    # TBD: to center the text
    for cell in iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(48)

def readCSVTo15Chunks_iter(filePath, lang="English"):
    '''Read csv file then get chars from api, and generate
     chars to random 15 chunks, finally yield the iter'''
    with open(filePath, 'r', encoding='utf-8') as f:
        csv_reader = reader(f)
        for row in csv_reader:
            chars, chars_len = getCharsFromAPI(row[1], lang)
            yield generateChars15Chunk(chars_len, chars)

def buildMultiSliderByCSV(ppt_name, chars_15chunks_iter):
    '''to build a ppt file with multi slider by
    reading each chars 15 chunks iter from csv file'''
    prs = Presentation()
    i = 1
    for chars in chars_15chunks_iter:
        addSlide(chars, prs, "SliderPlus {}".format(i))
        i += 1
    # create a unique file name (so that we don't have to open and close it during the testting)
    time_stamp = datetime.datetime.now().strftime('%Y%m%d%H%M%S%f')
    pptx_file_name = ppt_name + time_stamp + ".pptx"
    prs.save(pptx_file_name)
    return pptx_file_name

if __name__ == "__main__":
    chars_chunk_list = readCSVTo15Chunks_iter("quotes_english.csv")
    filename = buildMultiSliderByCSV("slider16_en_", chars_chunk_list)
    # chars_chunk_list = readCSVTo15Chunks_iter("quotes_telugu.csv", lang="Telugu")
    # filename = buildMultiSliderByCSV("slider16_telugu_", chars_chunk_list)
    print("Output PPT filename is: {}".format(filename))


