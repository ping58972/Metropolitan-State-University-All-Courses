
from pptx import Presentation
from pptx.util import Inches, Pt
import requests
import datetime
import json
import os

#========================= Calling WS API ==============
#  Each WS API will have a wrapper function in Python
# ======================================================
def getLogicalChars(input_str, language='English'):
    '''@string: input string for processing
       @language: language param,  default is English '''
    ws_api =  "https://indic-wp.thisisjava.com/api/getLogicalChars.php"
    params = {"string": input_str, "language": language}

    # On some servers, Apache may reject the crawlers.
    # To mimic an actual user, send this dummy header along
    # https://developer.mozilla.org/en-US/docs/Web/HTTP/Headers/User-Agent
    headers = {'User-Agent': ''}

    # get the response object
    response_obj = requests.get(url=ws_api, params=params, headers=headers)

    # get the json response
    json_response = response_obj.text

    # Get rid of UTF-8 Byte Order Mark (BOM)
    if json_response.startswith(u'\ufeff'):
        json_response = json_response.encode('utf8')[6:].decode('utf8')

    # Load the json response to convert it a dictionary
    json_dict = json.loads(json_response)

    # Get the logical characters (spaces are also counted)
    logical_characters = json_dict['data']

    # return the list
    return logical_characters

# ============= Process the logical characters =====
#  TBD: Overwrite this method with your own logic
# ======================================================
def processResponseData(logicalChars):
    # TBD
    return logicalChars


# ============= Create the PPTX presentation =====
# See https://python-pptx.readthedocs.io/en/latest/user/quickstart.html
#  You need to install python-pptx  (pip install python-pptx) in pycharm terminal
# ======================================================
def buildPresentation(processed_data):
    # innter function to return the cells of a table
    def iter_cells(table):
        for row in table.rows:
            for cell in row.cells:
                yield cell

    # constant for the slide layout
    # https://python-pptx.readthedocs.io/en/latest/user/table.html
    SLD_LAYOUT_TITLE_AND_CONTENT = 5

    # slide title
    title_txt = "SliderPlus"

    # create the presentation and set the slide using the slide layout
    pptx = Presentation()
    slide_layout = pptx.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
    slide = pptx.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_txt


    # cell dimensions for the table
    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)
    row_count, column_count = 4,4
    shape = slide.shapes.add_table(row_count, column_count, x, y, cx, cy)

    # set the table style
    # https://github.com/mbachtell/python-pptx-table-styles
    style_id = '{5940675A-B579-460E-94D1-54222C63F5DA}'
    tbl = shape._element.graphic.graphicData.tbl
    tbl[0][-1].text = style_id

    # populate the table with the content
    # TBD (this is dummy implementation)
    slider_table = shape.table
    for i in range(15):
        token = processed_data[i]
        cell = slider_table.cell(i // 4, i % 4)
        cell.text = token

    # set the font for the table cells
    # TBD: to center the text
    for cell in iter_cells(slider_table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(48)


    # create a unique file name (so that we don't have to open and close it during the testting)
    time_stamp = datetime.datetime.now().strftime('%Y%m%d%H%M%S%f')
    pptx_file_name = "slider16" + time_stamp + ".pptx"

    # save the presentation
    pptx.save(pptx_file_name)

    return pptx_file_name


# ============= Testign starts here =====
input_string = "Keep calm and program in python!"
logical_chars = getLogicalChars(input_string)
print("Input String :", input_string)
print("Logical Characters are: ", logical_chars)
processed_data = processResponseData(logical_chars)
pptx_file_name = buildPresentation(processed_data)

# view the presentation
# os.startfile(pptx_file_name)